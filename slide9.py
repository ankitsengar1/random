import pandas as pd
import numpy as np
from pptx.enum.dml import MSO_COLOR_TYPE
from pptx.enum.dml import MSO_THEME_COLOR_INDEX
from pptx.dml.color import RGBColor
import warnings
from formatting_functions import format_with_locale
import variables
from pptx.util import Pt,Inches
import matplotlib.pyplot as plt
from matplotlib.ticker import FuncFormatter
warnings.simplefilter(action='ignore', category=FutureWarning)

def edit_table(slide,df):
    for shape in slide.shapes:
        if shape.has_table:
            table = shape.table
            for col in range(0,len(table.columns)):
                for row in range(0,len(table.rows)):
                    for paragraph in table.cell(row,col).text_frame.paragraphs:
                        first_run = paragraph.runs[0] if paragraph.runs else paragraph.add_run()
                        font = first_run.font
                        font_name = font.name
                        font_size = font.size
                        font_bold = font.bold
                        font_italic = font.italic
                        font_underline = font.underline
                        text = paragraph.text
                        if font.color.type == MSO_COLOR_TYPE.SCHEME:
                            # Get RGB from theme color
                            theme_color = font.color.theme_color
                            font_rgb_flag = 1
                        elif font.color.type == None:
                            font_rgb_flag = 0
                        elif font.color.type == MSO_COLOR_TYPE.PRESET:
                            color_index = font.color.theme_color
                            font_rgb_flag = 3
                        else:
                            # Already RGB color, use it directly
                            rgb = font.color.rgb
                            font_rgb_flag = 2
                
                        # if row==0:
                        #     pass
                        # else:
                        paragraph.clear()      
                        new_run = paragraph.add_run()
                        new_run.font.name = font_name
                        new_run.font.size = font_size
                        new_run.font.bold = font_bold
                        new_run.font.underline = font_underline
                        new_run.font.italic = font_italic
                        # print(font_rgb_flag)
                        if font_rgb_flag == 1:
                            new_run.font.color.theme_color = theme_color
                        elif font_rgb_flag == 2:
                            new_run.font.color.rgb = rgb
                        elif font_rgb_flag == 3:
                            new_run.font.color.color_name = color_name = MSO_THEME_COLOR_INDEX(color_index).name

                        if row > 0: # leaving headers values
                            new_run.text = str(df.iloc[row-1,col])

                        elif row == 0: # headers
                            new_run.text = str(list(df.columns.values)[col])
                        
    return slide
#---------------------------------------------------------------------------------------------------------------------------------------------------------------------
def create_tagging_analysis_graph(slide,df):

    # Extract and transform the data
    df = df.set_index(" ").T  # Set first column as index and transpose

    # Extract values from DataFrame
    dates = df.index.values 
    total_bills = df["total_bills"].values
    loyalty_bills = df["loyalty_bills"].values

    x = np.arange(len(dates))  # Label locations
    width = 0.25  # Width of the bars
    bar_gap = 0.1  # Gap between bars

    # Creating the figure and axis
    fig, ax = plt.subplots(figsize=(6.02, 4.25))

    # Plot bars
    bars1 = ax.bar(x - (width/2 + bar_gap/2), total_bills, width, label='Total Bills', color='#4472c4')
    bars2 = ax.bar(x + (width/2 + bar_gap/2), loyalty_bills, width, label='Loyalty Bills', color='#ed7d31')

    # Add labels with a rectangle background
    offset = max(total_bills)

    for bar in bars1:
        ax.text(bar.get_x() + bar.get_width()/2, bar.get_height() + offset*0.01, format_with_locale(bar.get_height()), 
                ha='center', va='bottom', fontsize=10, color='black',
                bbox=dict(facecolor='white', edgecolor='white', boxstyle='round,pad=0.3'))

    for bar in bars2:
        ax.text(bar.get_x() + bar.get_width()/2, bar.get_height() - offset*0.1, format_with_locale(bar.get_height()), 
                ha='center', va='bottom', fontsize=10, color='black',
                bbox=dict(facecolor='white', edgecolor='white', boxstyle='round,pad=0.3'))

    # Customizing the plot
    ax.set_xlabel("", fontsize=12)
    ax.set_ylabel("", fontsize=12)
    ax.set_xticks(x)
    ax.set_xticklabels(dates, fontsize=10)
    #ax.legend(fontsize=11, loc='upper left')
    ax.legend(loc='upper center', bbox_to_anchor=(0.5, -0.1), ncol=2, frameon=False, fontsize=10)


    # ax.set_ylim(0, 12000)  # Adjusting y-axis limit
    # Dynamic y-axis limit with percentage padding
    max_value = max(max(total_bills), max(loyalty_bills))
    padding = max_value * 0.2  # 10% padding
    ax.set_ylim(0, max_value + padding)

    ax.spines['top'].set_visible(False)
    ax.spines['right'].set_visible(False)
    ax.spines['left'].set_visible(False)
    ax.xaxis.set_ticks_position('none')
    ax.yaxis.set_ticks_position('none')

    ax.yaxis.set_major_formatter(FuncFormatter(lambda x, _: format_with_locale(x)))

    # # Add horizontal grid lines
    ax.yaxis.grid(True, linestyle='-', linewidth=0.7, alpha=0.7, zorder=0)
    ax.set_axisbelow(True)


    # Show the plot
    plt.tight_layout()

    plt_path = variables.graph_path
    plt.savefig(plt_path,bbox_inches='tight')

    left = Inches(6.89)
    top = Inches(1.28)
    width = Inches(6.02)
    height = Inches(4.24)
    pic = slide.shapes.add_picture(plt_path, left, top, width, height)

    return slide

#---------------------------------------------------------------------------------------------------------------------------------------------------------------------
def create_df(chogori_ppt_xl_path,sheet_name):

    df = pd.read_excel(chogori_ppt_xl_path,sheet_name = sheet_name)
    df.rename(columns={df.columns[0]:" ",df.columns[1]:variables.prev_fourth_month,df.columns[2]:variables.prev_third_month,df.columns[3]:variables.prev_second_month,df.columns[4]:variables.prev_month,df.columns[5]:variables.current_month},inplace = True)
    df_g = df.copy()

    replacements = {'total_bills':'Total Bills', 'loyalty_bills': 'Loyalty Bills'}
    df[" "] = df[" "].replace(replacements)

    tagging = (df[df[" "] == "Loyalty Bills"].iloc[:, 1:].values / df[df[" "] == "Total Bills"].iloc[:, 1:].values)*100
    tagging_row = pd.DataFrame([["Tagging %"] + tagging.flatten().tolist()], columns=df.columns)

    df = pd.concat([df,tagging_row], ignore_index=True)

    target_columns = [df.columns[1], df.columns[2],df.columns[3],df.columns[4],df.columns[5]]

    for col in target_columns:
        df.loc[df[" "].isin(["Tagging %"]), col] = df.loc[df[" "].isin(["Tagging %"]), col].apply(lambda x: f"{round(x, 1)}%")

    for col in target_columns:
        df.loc[df[" "].isin(["Total Bills","Loyalty Bills"]), col] = df.loc[df[" "].isin(["Total Bills","Loyalty Bills"]), col].apply(format_with_locale)

    return df,df_g







