import pandas as pd
import numpy as np
from pptx.enum.dml import MSO_COLOR_TYPE
from pptx.enum.dml import MSO_THEME_COLOR_INDEX
from pptx.dml.color import RGBColor
import warnings
from formatting_functions import format_with_locale
import variables
warnings.simplefilter(action='ignore', category=FutureWarning)

def edit_table(slide,df_top,df_bottom):
    for i in [0,1]:
        if slide.shapes[i].has_table:
            table = slide.shapes[i].table
            for col in range(0,len(table.columns)):
                for row in range(0,len(table.rows)):
                    for paragraph in table.cell(row,col).text_frame.paragraphs:
                        first_run = paragraph.runs[0] if paragraph.runs else paragraph.add_run()
                        font = first_run.font
                        font_name = font.name
                        font_size = font.size
                        font_bold = font.bold
                        font_italic = font.italic
                        font_underline = font.underline
                        text = paragraph.text
                        if font.color.type == MSO_COLOR_TYPE.SCHEME:
                            # Get RGB from theme color
                            theme_color = font.color.theme_color
                            font_rgb_flag = 1
                        elif font.color.type == None:
                            font_rgb_flag = 0
                        elif font.color.type == MSO_COLOR_TYPE.PRESET:
                            color_index = font.color.theme_color
                            font_rgb_flag = 3
                        else:
                            # Already RGB color, use it directly
                            rgb = font.color.rgb
                            font_rgb_flag = 2

                        # if row==0:
                        #     pass
                        # else:
                        paragraph.clear()                    
                        new_run = paragraph.add_run()
                        new_run.font.name = font_name
                        new_run.font.size = font_size
                        new_run.font.bold = font_bold
                        new_run.font.underline = font_underline
                        new_run.font.italic = font_italic
                        # print(font_rgb_flag)
                        if font_rgb_flag == 1:
                            new_run.font.color.theme_color = theme_color
                        elif font_rgb_flag == 2:
                            new_run.font.color.rgb = rgb
                        elif font_rgb_flag == 3:
                            new_run.font.color.color_name = color_name = MSO_THEME_COLOR_INDEX(color_index).name

                        if i==0:
                            if row > 0: # leaving headers values
                                new_run.text = str(df_top.iloc[row-1,col])

                            elif row == 0: # headers
                                new_run.text = str(list(df_top.columns.values)[col])

                        elif i==1:
                            if row > 0: # leaving headers values
                                new_run.text = str(df_bottom.iloc[row-1,col])

                            elif row == 0: # headers
                                new_run.text = str(list(df_bottom.columns.values)[col])

    return slide



def create_df(chogori_ppt_xl_path,sheet_name):

    df = pd.read_excel(chogori_ppt_xl_path,sheet_name = sheet_name)
    df["total_bills"] = df["bills"] + df["nonloyalty_bills"]
    df["Tagging%"] = df["bills"]/df["total_bills"]*100
    df.sort_values(by="Tagging%", ascending=False,inplace = True)
    df = df[["lpaas_store","Tagging%"]]

    df_top = df.head(10)
    df_top.rename(columns = {"lpaas_store":"TOP 10 TAGGING STORES"},inplace = True)
    df_top["Tagging%"] = df_top["Tagging%"].round(0).astype(int).astype(str)+"%"
    df_bottom = df.tail(10)
    df_bottom.sort_values(by="Tagging%", ascending=True,inplace = True)
    df_bottom.rename(columns = {"lpaas_store":"BOTTOM 10 TAGGING STORES"},inplace = True)
    df_bottom["Tagging%"] = df_bottom["Tagging%"].round(0).astype(int).astype(str)+"%"


    return df_top,df_bottom







