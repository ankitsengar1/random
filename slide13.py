import pandas as pd
import numpy as np
from pptx.enum.dml import MSO_COLOR_TYPE
from pptx.enum.dml import MSO_THEME_COLOR_INDEX
from pptx.dml.color import RGBColor
import warnings
from formatting_functions import format_with_locale
import variables
from pptx.util import Pt,Inches
import matplotlib.pyplot as plt
from matplotlib.ticker import FuncFormatter
warnings.simplefilter(action='ignore', category=FutureWarning)

def edit_table(slide,df_1,df_2):
    for i in [1,2]:
        if slide.shapes[i].has_table:
            table = slide.shapes[i].table
            for col in range(0,len(table.columns)):
                for row in range(0,len(table.rows)):
                    for paragraph in table.cell(row,col).text_frame.paragraphs:
                        first_run = paragraph.runs[0] if paragraph.runs else paragraph.add_run()
                        font = first_run.font
                        font_name = font.name
                        font_size = font.size
                        font_bold = font.bold
                        font_italic = font.italic
                        font_underline = font.underline
                        text = paragraph.text
                        if font.color.type == MSO_COLOR_TYPE.SCHEME:
                            # Get RGB from theme color
                            theme_color = font.color.theme_color
                            font_rgb_flag = 1
                        elif font.color.type == None:
                            font_rgb_flag = 0
                        elif font.color.type == MSO_COLOR_TYPE.PRESET:
                            color_index = font.color.theme_color
                            font_rgb_flag = 3
                        else:
                            # Already RGB color, use it directly
                            rgb = font.color.rgb
                            font_rgb_flag = 2

                        # if (row>=0 and col==0):
                        #     pass
                        # else:
                        paragraph.clear()                    
                        new_run = paragraph.add_run()
                        new_run.font.name = font_name
                        new_run.font.size = font_size
                        new_run.font.bold = font_bold
                        new_run.font.underline = font_underline
                        new_run.font.italic = font_italic
                        # print(font_rgb_flag)
                        if font_rgb_flag == 1:
                            new_run.font.color.theme_color = theme_color
                        elif font_rgb_flag == 2:
                            new_run.font.color.rgb = rgb
                        elif font_rgb_flag == 3:
                            new_run.font.color.color_name = color_name = MSO_THEME_COLOR_INDEX(color_index).name

                        if i==1:
                            if row > 0: # leaving headers values
                                new_run.text = str(df_1.iloc[row-1,col])

                            elif row == 0: # headers
                                new_run.text = str(list(df_1.columns.values)[col])

                        elif i==2:
                            if row > 0: # leaving headers values
                                new_run.text = str(df_2.iloc[row-1,col])

                            elif row == 0: # headers
                                new_run.text = str(list(df_2.columns.values)[col])

    return slide
#-------------------------------------------------------------------------------------------------------------------------------------------------------------------------
def create_overall_program_graph(slide,df):

    # Extract values from DataFrame
    time = df["month_tag"].values
    onetimer = df["one_timer_count"].values
    repeater = df["repeater_count"].values

    x = np.arange(len(time))  # Label locations
    width = 0.25  # Width of the bars
    bar_gap = 0.1  # Gap between bars

    # Creating the figure and axis
    fig, ax = plt.subplots(figsize=(6.43, 2.88))

    # Plot bars
    bars1 = ax.bar(x - (width/2 + bar_gap/2), onetimer, width, label='One Timer', color='#4472c4')
    bars2 = ax.bar(x + (width/2 + bar_gap/2), repeater, width, label='Repeater', color='#ed7d31')

    ## Add labels with a rectangle background
    # offset = max(onetimer)

    # for bar in bars1:
    #     ax.text(bar.get_x() + bar.get_width()/2, bar.get_height() + offset*0.01, format_with_locale(bar.get_height()), 
    #             ha='center', va='bottom', fontsize=10, color='black',
    #             bbox=dict(facecolor='white', edgecolor='white', boxstyle='round,pad=0.3'))

    # for bar in bars2:
    #     ax.text(bar.get_x() + bar.get_width()/2, bar.get_height() - offset*0.1, format_with_locale(bar.get_height()), 
    #             ha='center', va='bottom', fontsize=10, color='black',
    #             bbox=dict(facecolor='white', edgecolor='white', boxstyle='round,pad=0.3'))

    # Customizing the plot
    ax.set_xlabel("", fontsize=12)
    ax.set_ylabel("", fontsize=12)
    ax.set_xticks(x)
    ax.set_xticklabels(time, fontsize=10)
    #ax.legend(fontsize=11, loc='upper left')
    ax.legend(loc='upper center', bbox_to_anchor=(0.5, -0.1), ncol=2, frameon=False, fontsize=10)

    # Dynamic y-axis limit with percentage padding
    max_value = max(max(onetimer), max(repeater))
    padding = max_value * 0.2  # 10% padding
    ax.set_ylim(0, max_value + padding)

    ax.spines['top'].set_visible(False)
    ax.spines['right'].set_visible(False)
    ax.spines['left'].set_visible(False)
    ax.xaxis.set_ticks_position('none')
    ax.yaxis.set_ticks_position('none')

    ax.yaxis.set_major_formatter(FuncFormatter(lambda x, _: format_with_locale(x)))

    # # Add horizontal grid lines
    ax.yaxis.grid(True, linestyle='-', linewidth=0.7, alpha=0.7, zorder=0)
    ax.set_axisbelow(True)


    # Show the plot
    plt.tight_layout()

    plt_path = variables.graph_path
    plt.savefig(plt_path,bbox_inches='tight')

    left = Inches(6.37)
    top = Inches(0.94)
    width = Inches(6.42)
    height = Inches(2.88)
    pic = slide.shapes.add_picture(plt_path, left, top, width, height)

    return slide

def create_df(chogori_ppt_xl_path,sheet_name_1,sheet_name_2,sheet_name_3,overall_abv):

    df = pd.read_excel(chogori_ppt_xl_path,sheet_name = sheet_name_1)
    df_1 =df[["kpis","prev_month_onetimer","prev_month_repeater"]]
    df_2 = df[["kpis","curr_month_onetimer","curr_month_repeater"]]
    df_1.rename(columns = {"kpis":variables.prev_month,"prev_month_onetimer":"One Timer","prev_month_repeater":"Repeat"},inplace = True)
    df_2.rename(columns = {"kpis":variables.current_month,"curr_month_onetimer":"One Timer","curr_month_repeater":"Repeat"},inplace = True)

    target_columns = ["One Timer","Repeat"]
    cr_list = ["sales"]

    # for previous month
    for col in target_columns:
        df_1.loc[~df_1[df_1.columns[0]].isin(cr_list), col] = df_1.loc[~df_1[df_1.columns[0]].isin(cr_list), col].apply(format_with_locale)

    for col in target_columns:
        df_1.loc[df_1[df_1.columns[0]].isin(cr_list), col] = df_1.loc[df_1[df_1.columns[0]].isin(cr_list), col].apply(
            lambda x: f"{round(x / 10000000, 1)} Cr"
        )
    # for current month
    for col in target_columns:
        df_2.loc[~df_2[df_2.columns[0]].isin(cr_list), col] = df_2.loc[~df_2[df_2.columns[0]].isin(cr_list), col].apply(format_with_locale)

    for col in target_columns:
        df_2.loc[df_2[df_2.columns[0]].isin(cr_list), col] = df_2.loc[df_2[df_2.columns[0]].isin(cr_list), col].apply(
            lambda x: f"{round(x / 10000000, 1)} Cr"
        )

    replacements = {'customers':'Customers', 'sales': 'Sales',"bills":"Bills","atv":"ATV","amv":"AMV"}
    df_1[df_1.columns[0]] = df_1[df_1.columns[0]].replace(replacements)
    df_2[df_2.columns[0]] = df_2[df_2.columns[0]].replace(replacements)

    #Data for graph
    df_g = pd.read_excel(chogori_ppt_xl_path,sheet_name = sheet_name_2)

    # data for text
    one_timer_customer = df_2.loc[df_2.iloc[:, 0] == "Customers", "One Timer"].values[0]
    repeater_customer = df_2.loc[df_2.iloc[:, 0] == "Customers", "Repeat"].values[0]
    abv_repeat = df_2.loc[df_2.iloc[:, 0] == "ATV", "Repeat"].values[0]

    df_t = pd.read_excel(chogori_ppt_xl_path,sheet_name = sheet_name_3)

    # for current month
    c_max = str(round(df_t["monthly_sales"].max() / 10000000, 1)) + "Cr"
    c_min = str(round(df_t["monthly_sales"].min() / 10000000, 1)) + "Cr"
    c_day_max = df_t.loc[df_t["monthly_sales"].idxmax(), "day"]
    c_day_min = df_t.loc[df_t["monthly_sales"].idxmin(), "day"]

    # for finantial year
    fy_max = str(round(df_t["fy_sales"].max() / 10000000, 1)) + "Cr"
    fy_min = str(round(df_t["fy_sales"].min() / 10000000, 1)) + "Cr"
    fy_day_max = df_t.loc[df_t["fy_sales"].idxmax(), "day"]
    fy_day_min = df_t.loc[df_t["fy_sales"].idxmin(), "day"]

    old_text = ["One Timer : (curr_month) - onetimer",
                "Repeat customers : (curr_month) – repeater",
                "Average bill value of Loyal Members curr_month: overall_abv",
                "Average bill value of Repeat Members curr_month: repeat_abv",
                "curr_month- : c_day_max: Max sale – c_max, c_day_min : Min Sale – c_min",
                "FY(fy) fy_day_max: Max sale – fy_max, fy_day_min: Min Sale – fy_min"

    ]
    new_text = [f"One Timer : ({variables.current_month}) - {one_timer_customer}",
                f"Repeat customers : ({variables.current_month}) - {repeater_customer}",
                f"Average bill value of Loyal Members {variables.current_month} : {overall_abv}",
                f"Average bill value of Repeat Members {variables.current_month} : {abv_repeat}",
                f"{variables.current_month} - : {c_day_max} : Max sale - {c_max}, {c_day_min} : Min Sale - {c_min}",
                f"FY({variables.fy}) - : {fy_day_max} : Max sale - {fy_max}, {fy_day_min} : Min Sale - {fy_min}"
       
    ]

    return df_1,df_2,df_g,old_text,new_text