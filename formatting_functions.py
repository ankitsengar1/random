import pandas as pd
import numpy as np
import locale
from pptx.enum.dml import MSO_COLOR_TYPE,MSO_THEME_COLOR_INDEX
from pptx.dml.color import RGBColor
from matplotlib.colors import LinearSegmentedColormap
# from pptx.util import Pt,Inches
import matplotlib.pyplot as plt
import math

def format_with_locale(value):
    try:
        locale.setlocale(locale.LC_NUMERIC, 'en_IN')
        # Round to the nearest integer and apply Indian locale formatting
        rounded_value = round(value)
        return locale.format_string("%d", rounded_value, grouping=True)
    except Exception as e:
        return value  
    
#---------------------------------------------------------------------------------------------------------------------------------------------------------------
def change_text_box(slide_1,old_text,new_text):
    # print(old_text,new_text)
    for shape in slide_1.shapes:
        if shape.has_text_frame:
            for paragraph in shape.text_frame.paragraphs:
                # print(str(paragraph.text).replace("\xa0", " ").strip())
                # print(old_text)
                # print('-----------')
                if str(paragraph.text).replace("\xa0", " ").strip() == old_text:
                    first_run = paragraph.runs[0] if paragraph.runs else paragraph.add_run()
                    font = first_run.font
                    font_name = font.name
                    font_size = font.size
                    font_bold = font.bold
                    font_italic = font.italic
                    font_underline = font.underline
                    if font.color.type == MSO_COLOR_TYPE.SCHEME:
                        theme_color = font.color.theme_color
                        font_rgb_flag = 1
                    elif font.color.type == None:
                        font_rgb_flag = 0
                    elif font.color.type == MSO_COLOR_TYPE.PRESET:
                        color_index = font.color.theme_color
                        font_rgb_flag = 3
                    else:
                        rgb = font.color.rgb
                        font_rgb_flag = 2
            
                    paragraph.clear()
                    new_run = paragraph.add_run()
                    new_run.text = new_text
                    new_run.font.name = font_name
                    new_run.font.size = font_size
                    new_run.font.bold = font_bold
                    new_run.font.underline = font_underline
                    new_run.font.italic = font_italic
                    if font_rgb_flag == 1:
                        new_run.font.color.theme_color = theme_color
                    elif font_rgb_flag == 2:
                        new_run.font.color.rgb = rgb
                    elif font_rgb_flag == 3:
                        new_run.font.color.color_name = MSO_THEME_COLOR_INDEX(color_index).name
#--------------------------------------------------------------------------------------------------------------------------------------------------------------
#code to check index position of table
# table_indices = []
# for index, shape in enumerate(slide.shapes):
#     if shape.has_table:
#         table_indices.append(index)
# print("######################################################")
# print(table_indices)
# print("#####################################################")
# for i in [3,4]:  # Modify this if you want to loop through other table indices
#     if slide.shapes[i].has_table:
#         table = slide.shapes[i].table
#         print(f"Number of columns: {len(table.columns)}")
#         print(f"Number of rows: {len(table.rows)}")

#         for col in range(len(table.columns)):
#             for row in range(len(table.rows)):
#                 # Get the text content of the current cell
#                 cell_text = table.cell(row, col).text_frame.text if table.cell(row, col).text_frame else ""
                
#                 # Print row number, column number, and cell content
#                 print(f"Cell[{row}, {col}] Content: {cell_text}")