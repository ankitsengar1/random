import pandas as pd
import numpy as np
from pptx.enum.dml import MSO_COLOR_TYPE
from pptx.enum.dml import MSO_THEME_COLOR_INDEX
from pptx.dml.color import RGBColor
import warnings
from formatting_functions import format_with_locale
import variables
warnings.simplefilter(action='ignore', category=FutureWarning)

def edit_table(slide,df_1,df_2):
    for i in [0,1]:
        if slide.shapes[i].has_table:
            table = slide.shapes[i].table
            for col in range(0,len(table.columns)):
                for row in range(0,len(table.rows)):
                    for paragraph in table.cell(row,col).text_frame.paragraphs:
                        first_run = paragraph.runs[0] if paragraph.runs else paragraph.add_run()
                        font = first_run.font
                        font_name = font.name
                        font_size = font.size
                        font_bold = font.bold
                        font_italic = font.italic
                        font_underline = font.underline
                        text = paragraph.text
                        if font.color.type == MSO_COLOR_TYPE.SCHEME:
                            # Get RGB from theme color
                            theme_color = font.color.theme_color
                            font_rgb_flag = 1
                        elif font.color.type == None:
                            font_rgb_flag = 0
                        elif font.color.type == MSO_COLOR_TYPE.PRESET:
                            color_index = font.color.theme_color
                            font_rgb_flag = 3
                        else:
                            # Already RGB color, use it directly
                            rgb = font.color.rgb
                            font_rgb_flag = 2

                        # if row==0:
                        #     pass
                        # else:
                        paragraph.clear()                    
                        new_run = paragraph.add_run()
                        new_run.font.name = font_name
                        new_run.font.size = font_size
                        new_run.font.bold = font_bold
                        new_run.font.underline = font_underline
                        new_run.font.italic = font_italic
                        # print(font_rgb_flag)
                        if font_rgb_flag == 1:
                            new_run.font.color.theme_color = theme_color
                        elif font_rgb_flag == 2:
                            new_run.font.color.rgb = rgb
                        elif font_rgb_flag == 3:
                            new_run.font.color.color_name = color_name = MSO_THEME_COLOR_INDEX(color_index).name

                        if i==1:
                            if row > 0: # leaving headers values
                                new_run.text = str(df_1.iloc[row-1,col])

                            elif row == 0: # headers
                                new_run.text = str(list(df_1.columns.values)[col])

                        elif i==0:
                            if row > 0: # leaving headers values
                                new_run.text = str(df_2.iloc[row-1,col])

                            elif row == 0: # headers
                                new_run.text = str(list(df_2.columns.values)[col])

    return slide

def create_df(chogori_ppt_xl_path,sheet_name_1,sheet_name_2):

    df_1 = pd.read_excel(chogori_ppt_xl_path,sheet_name = sheet_name_1)
    df_1.rename(columns = {"year":"Year","points_collected":"Points Earned","points_reedemed":"Points Burn"},inplace = True)
    df_1["Burn/Earn Ratio"] = df_1["Points Burn"]/df_1["Points Earned"]*100
    df_1.fillna(0,inplace = True)
    df_1["Burn/Earn Ratio"] = df_1["Burn/Earn Ratio"].round(1).astype(str)+"%"
    df_1["Points Earned"] = (df_1["Points Earned"]/100000).round(2).astype(str)+ "L"
    df_1["Points Burn"] = (df_1["Points Burn"]/100000).round(2).astype(str)+ "L"
    df_1.replace(["0.0L","0.0%"],"-",inplace = True)

    # creating second dataframe
    df_2 = pd.read_excel(chogori_ppt_xl_path,sheet_name = sheet_name_2)
    df_2.rename(columns = {"mom":"Month","points_collected":"Points Earned","points_reedemed":"Points Burn"},inplace = True)
    df_2["Burn/Earn Ratio"] = df_2["Points Burn"]/df_2["Points Earned"]*100
    df_2.fillna(0,inplace = True)
    df_2["Burn/Earn Ratio"] = df_2["Burn/Earn Ratio"].round(1).astype(str)+"%"
    df_2["Points Earned"] = (df_2["Points Earned"]/100000).round(2).astype(str)+ "L"
    df_2["Points Burn"] = (df_2["Points Burn"]/100000).round(2).astype(str)+ "L"

    replacements = {'prev_third_month':variables.prev_third_month, 'prev_second_month': variables.prev_second_month,'prev_month': variables.prev_month,"curr_month":variables.current_month}
    df_2["Month"] = df_2["Month"].replace(replacements)

    return df_1,df_2







