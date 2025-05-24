import pandas as pd
import numpy as np
from pptx.enum.dml import MSO_COLOR_TYPE
from pptx.enum.dml import MSO_THEME_COLOR_INDEX
from pptx.dml.color import RGBColor
import warnings
from formatting_functions import format_with_locale
import variables
warnings.simplefilter(action='ignore', category=FutureWarning)

def edit_table(slide,df_1,df_2):
    for i in [1,2]:
        if slide.shapes[i].has_table:
            table = slide.shapes[i].table
            for col in range(0,len(table.columns)):
                for row in range(0,len(table.rows)):
                    for paragraph in table.cell(row,col).text_frame.paragraphs:
                        first_run = paragraph.runs[0] if paragraph.runs else paragraph.add_run()
                        font = first_run.font
                        font_name = font.name
                        font_size = font.size
                        font_bold = font.bold
                        font_italic = font.italic
                        font_underline = font.underline
                        text = paragraph.text
                        if font.color.type == MSO_COLOR_TYPE.SCHEME:
                            # Get RGB from theme color
                            theme_color = font.color.theme_color
                            font_rgb_flag = 1
                        elif font.color.type == None:
                            font_rgb_flag = 0
                        elif font.color.type == MSO_COLOR_TYPE.PRESET:
                            color_index = font.color.theme_color
                            font_rgb_flag = 3
                        else:
                            # Already RGB color, use it directly
                            rgb = font.color.rgb
                            font_rgb_flag = 2

                        if i==2:
                            if (row==0 and col==0)or(row==1 and col>=0):
                                pass

                            else:
                                paragraph.clear()                    
                                new_run = paragraph.add_run()
                                new_run.font.name = font_name
                                new_run.font.size = font_size
                                new_run.font.bold = font_bold
                                new_run.font.underline = font_underline
                                new_run.font.italic = font_italic
                                # print(font_rgb_flag)
                                if font_rgb_flag == 1:
                                    new_run.font.color.theme_color = theme_color
                                elif font_rgb_flag == 2:
                                    new_run.font.color.rgb = rgb
                                elif font_rgb_flag == 3:
                                    new_run.font.color.color_name = color_name = MSO_THEME_COLOR_INDEX(color_index).name

                                if i==2:
                                    if row > 1:
                                        new_run.text = str(df_1.iloc[row-2,col])
                                    elif row == 0 and col==1:
                                        new_run.text = variables.prev_month
                                    elif row == 0 and col==4:
                                         new_run.text = variables.current_month

                        else:
                            paragraph.clear()                    
                            new_run = paragraph.add_run()
                            new_run.font.name = font_name
                            new_run.font.size = font_size
                            new_run.font.bold = font_bold
                            new_run.font.underline = font_underline
                            new_run.font.italic = font_italic
                            # print(font_rgb_flag)
                            if font_rgb_flag == 1:
                                new_run.font.color.theme_color = theme_color
                            elif font_rgb_flag == 2:
                                new_run.font.color.rgb = rgb
                            elif font_rgb_flag == 3:
                                new_run.font.color.color_name = color_name = MSO_THEME_COLOR_INDEX(color_index).name

                            if i==1:
                                if row > 0:
                                    new_run.text = str(df_2.iloc[row-1,col])
                                elif row == 0:
                                    new_run.text = str(list(df_2.columns.values)[col])

                
    return slide

def create_df(chogori_ppt_xl_path,sheet_name_1,sheet_name_2):

    df1 = pd.read_excel(chogori_ppt_xl_path,sheet_name = sheet_name_1)
    df1.rename(columns={df1.columns[0]:"KPIS"},inplace = True)
    replacements = {'customer':'Total Customers', 'points_collected': 'Transaction Points Collected','points_reedemed':'Total Points Redeemed','redeemers':'Point Redeemers',
                   'redemption_bills':'Redemption Bills','redemption_sales':'Redemption Sales','accrued_customer':'Accrued customer','points_issued':'Bonus Points Issued'
                    }
    df1["KPIS"] = df1["KPIS"].replace(replacements)
    df1.insert(3,"overall_prev_month",(df1["offline_prev_month"]+df1["online_prev_month"]))
    df1.insert(6,"overall_curr_month",(df1["offline_curr_month"]+df1["online_curr_month"]))
    # print(df1)

    points_redeemers = (df1[df1["KPIS"] == "Point Redeemers"].iloc[:, 1:].values / df1[df1["KPIS"] == "Total Customers"].iloc[:, 1:].values)*100
    points_redeemers_row = pd.DataFrame([["Point Redeemers %"] + points_redeemers.flatten().tolist()], columns=df1.columns)

    redemption_rate = (df1[df1["KPIS"] == "Total Points Redeemed"].iloc[:, 1:].values / df1[df1["KPIS"] == "Transaction Points Collected"].iloc[:, 1:].values)*100
    redemptiom_rate_row = pd.DataFrame([["Point Redemption Rate"] + redemption_rate.flatten().tolist()], columns=df1.columns)

    total_points = (df1[df1["KPIS"] == "Transaction Points Collected"].iloc[:, 1:].values + df1[df1["KPIS"] == "Bonus Points Issued"].iloc[:, 1:].values)
    total_points_row = pd.DataFrame([["Total Points Issued"] + total_points.flatten().tolist()], columns=df1.columns)

    df1 = pd.concat([df1,points_redeemers_row,redemptiom_rate_row,total_points_row], ignore_index=True)
    df1.fillna(0,inplace = True)

    percent_kpis = ["Point Redeemers %", "Point Redemption Rate"]

    exceptions = ["Point Redeemers %", "Point Redemption Rate","Transaction Points Collected","Total Points Redeemed","Redemption Sales","Total Points Issued"]

    target_per = ["Point Redeemers %", "Point Redemption Rate"]
    target_kpis_lac = ["Transaction Points Collected","Total Points Redeemed","Redemption Sales","Total Points Issued"]

    target_columns = [df1.columns[1], df1.columns[3],df1.columns[4],df1.columns[6]]

    for col in target_columns:
        df1.loc[df1["KPIS"].isin(target_per), col] = df1.loc[df1["KPIS"].isin(target_per), col].apply(lambda x: f"{round(x, 1)}%")

    for col in target_columns:
        df1.loc[~df1["KPIS"].isin(exceptions), col] = df1.loc[~df1["KPIS"].isin(exceptions), col].apply(format_with_locale)

    # for col in target_columns:
    #     df.loc[df["KPIS"].isin(target_kpis_inr_cr), col] = df.loc[df["KPIS"].isin(target_kpis_inr_cr), col].apply(
    #         lambda x: f"{round(x / 10000000, 1)} Cr"
    #     )
    for col in target_columns:
        df1.loc[df1["KPIS"].isin(target_kpis_lac), col] = df1.loc[df1["KPIS"].isin(target_kpis_lac), col].apply(
            lambda x: f"{round(x / 100000, 2)} L"
        )
        
    desired_order = ["Total Customers","Transaction Points Collected","Total Points Redeemed","Point Redeemers","Point Redeemers %","Redemption Bills",
                     "Redemption Sales","Accrued customer","Bonus Points Issued","Total Points Issued","Point Redemption Rate"
        ]
    df1 = df1.set_index("KPIS") 
    df1 = df1.loc[desired_order].reset_index()
    df1.replace(0,"-",inplace = True)

    # creating second dataframe ---------------------------------------------------------------------------------------------------------------------------------
    df2 = pd.read_excel(chogori_ppt_xl_path,sheet_name = sheet_name_2)
    df2.rename(columns={df2.columns[0]:"Month",df2.columns[1]:variables.prev_month,df2.columns[2]:variables.current_month},inplace = True)
    replacements = {'issued':'ISSUED COUPONS', 'coupons_redeemed': 'REDEEMED COUPONS','redeemers':'COUPON REDEEMERS',"discount":"DISCOUNT VALUE"
                    }
    df2["Month"] = df2["Month"].replace(replacements)

    coupon_rrate = (df2[df2["Month"] == "REDEEMED COUPONS"].iloc[:, 1:].values / df2[df2["Month"] == "ISSUED COUPONS"].iloc[:, 1:].values)*100
    coupon_rrate_row = pd.DataFrame([["COUPON REDMPTION RATE"] + coupon_rrate.flatten().tolist()], columns=df2.columns)

    df2 = pd.concat([df2,coupon_rrate_row], ignore_index=True)

    target_columns = [df2.columns[1], df2.columns[2]]

    for col in target_columns:
        df2.loc[df2["Month"].isin(['COUPON REDMPTION RATE']), col] = df2.loc[df2["Month"].isin(['COUPON REDMPTION RATE']), col].apply(lambda x: f"{round(x, 1)}%")

    for col in target_columns:
        df2.loc[~df2["Month"].isin(["COUPON REDMPTION RATE","DISCOUNT VALUE"]), col] = df2.loc[~df2["Month"].isin(["COUPON REDMPTION RATE","DISCOUNT VALUE"]), col].apply(format_with_locale)

    # # for col in target_columns:
    # #     df.loc[df["KPIS"].isin(target_kpis_inr_cr), col] = df.loc[df["KPIS"].isin(target_kpis_inr_cr), col].apply(
    # #         lambda x: f"{round(x / 10000000, 1)} Cr"
    # #     )
    for col in target_columns:
        df2.loc[df2["Month"].isin(['DISCOUNT VALUE']), col] = df2.loc[df2["Month"].isin(['DISCOUNT VALUE']), col].apply(
            lambda x: f"{round(x / 100000, 2)} L"
        )
    desired_order = ["ISSUED COUPONS","REDEEMED COUPONS","COUPON REDEEMERS","DISCOUNT VALUE","COUPON REDMPTION RATE"]

    return df1,df2







