from pptx import Presentation
import variables,formatting_functions
import slide5,slide6,slide7,slide8,slide9,slide10,slide11,slide12,slide13,slide15,slide16,slide14

# chogori excel path
chogori_ppt_xl_path = variables.ppt_xl_path
# Create a presentation object
prs = Presentation(variables.template_path)

#-----------------------------------------------------------------------------------------------------------------------------------------------
#Slide 1
slide_num = 1
slide_1 = prs.slides[slide_num-1]

old_text =[
    "CRM MONTHLY REVIEW",
    "curr_month"
]

new_text = [
    "CRM MONTHLY REVIEW",
    f"{variables.curr_month_fname}"]

for i in range(len(old_text)):
    formatting_functions.change_text_box(slide_1,old_text[i],new_text[i])
#-------------------------------------------------------------------------------------------------------------------------------------------------
#Slide 2
slide_num = 2
slide_2 = prs.slides[slide_num-1]

old_text =[
    "Detailed KPI View Overall – curr_month",
    "ABV Band Wise Distribution Of Customer, Sales And Bills In curr_month",
    "Top/Bottom 10 Stores Based On Bill Tagging – curr_month",
    "Top/Bottom 10 Stores Based on Redemption – curr_month"
]

new_text = [
    f"Detailed KPI View Overall - {variables.current_month}",
    f"ABV Band Wise Distribution Of Customer, Sales And Bills In {variables.current_month}",
    f"Top/Bottom 10 Stores Based On Bill Tagging - {variables.current_month}",
    f"Top/Bottom 10 Stores Based on Redemption - {variables.current_month}",

    ]

for i in range(len(old_text)):
    formatting_functions.change_text_box(slide_2,old_text[i],new_text[i])
#-------------------------------------------------------------------------------------------------------------------------------------------------
# Slide 5
slide_num = 5
slide_5 = prs.slides[slide_num-1]
sheet_name = "overall_kpis"
df,overall_abv = slide5.create_df(chogori_ppt_xl_path,sheet_name)
slide_5 = slide5.edit_table(slide_5,df)

old_text = "DETAILED KPI VIEW FOR curr_month – OVERALL"
new_text = f"DETAILED KPI VIEW FOR {variables.current_month} - OVERALL"
slide_5 = formatting_functions.change_text_box(slide_5,old_text,new_text)
#-------------------------------------------------------------------------------------------------------------------------------------------------
# Slide 6
slide_num = 6
slide_6 = prs.slides[slide_num-1]
sheet_name = "atv_band"
df = slide6.create_df(chogori_ppt_xl_path,sheet_name)
slide_6 = slide6.edit_table(slide_6,df)

old_text = "ABV BAND WISE DISTRIBUTION OF CUSTOMER, SALES AND BILLS IN curr_month"
new_text = f"ABV BAND WISE DISTRIBUTION OF CUSTOMER, SALES AND BILLS IN {variables.current_month}"
slide_6 = formatting_functions.change_text_box(slide_6,old_text,new_text)
#-------------------------------------------------------------------------------------------------------------------------------------------------
# Slide 7
slide_num = 7
slide_7 = prs.slides[slide_num-1]
sheet_name_1 = "points_data"
sheet_name_2 = "coupon_data"
df_1,df_2 = slide7.create_df(chogori_ppt_xl_path,sheet_name_1,sheet_name_2)
slide_7 = slide7.edit_table(slide_7,df_1,df_2)

old_text = "POINT & COUPON ANALYSIS – ISSUANCE & REDEMPTION prev_month & curr_month"
new_text = f"POINT & COUPON ANALYSIS – ISSUANCE & REDEMPTION {variables.prev_month} & {variables.current_month}"
slide_7 = formatting_functions.change_text_box(slide_7,old_text,new_text)
#---------------------------------------------------------------------------------------------------------------------------------------------
# Slide 8
slide_num = 8
slide_8 = prs.slides[slide_num-1]
sheet_name = "points_distribution"
df = slide8.create_df(chogori_ppt_xl_path,sheet_name)
slide_8 = slide8.edit_table(slide_8,df)
#-------------------------------------------------------------------------------------------------------------------------------------------------
# Slide 9
slide_num = 9
slide_9 = prs.slides[slide_num-1]
sheet_name = "tagging_analysis_mom"
df,df_g = slide9.create_df(chogori_ppt_xl_path,sheet_name)
slide_9 = slide9.edit_table(slide_9,df)
slide_9 = slide9.create_tagging_analysis_graph(slide_9,df_g)
#-------------------------------------------------------------------------------------------------------------------------------------------------
# Slide 10
slide_num = 10
slide_10 = prs.slides[slide_num-1]
sheet_name = "loyalty_data_storewise"
df_top,df_bottom = slide10.create_df(chogori_ppt_xl_path,sheet_name)
slide_10 = slide10.edit_table(slide_10,df_top,df_bottom)

old_text = "BOTTOM/TOP STORES - BASED ON TAGGING – (curr_month)"
new_text = f"BOTTOM/TOP STORES - BASED ON TAGGING - ({variables.current_month})"
slide_10 = formatting_functions.change_text_box(slide_10,old_text,new_text)
#-------------------------------------------------------------------------------------------------------------------------------------------------
# Slide 11
slide_num = 11
slide_11 = prs.slides[slide_num-1]
sheet_name_1 = "burn_to_earn_yrwise"
sheet_name_2 = "burn_to_earn_monthwise"
df_1,df_2 = slide11.create_df(chogori_ppt_xl_path,sheet_name_1,sheet_name_2)
slide_11 = slide11.edit_table(slide_11,df_1,df_2)
#-------------------------------------------------------------------------------------------------------------------------------------------------
# Slide 12
slide_num = 12
slide_12 = prs.slides[slide_num-1]
sheet_name = "storewise_redeemption"
df_1,df_2 = slide12.create_df(chogori_ppt_xl_path,sheet_name)
slide_12 = slide12.edit_table(slide_12,df_1,df_2)

old_text =[
    "BOTTOM/TOP STORES - BASED ON REDEMPTION – (curr_month)",
    "Stores with Zero/Low Redemption (curr_month)",
    "Stores with Maximum Redemption (curr_month)"
]

new_text = [
    f"BOTTOM/TOP STORES - BASED ON REDEMPTION - ({variables.current_month})",
    f"Stores with Zero/Low Redemption ({variables.current_month})",
    f"Stores with Maximum Redemption ({variables.current_month})"
    ]

for i in range(len(old_text)):
    formatting_functions.change_text_box(slide_12,old_text[i],new_text[i])
#-------------------------------------------------------------------------------------------------------------------------------------------------
# Slide 13
slide_num = 13
slide_13 = prs.slides[slide_num-1]
sheet_name_1 = "overall_program_kpis"
sheet_name_2 = "repeater_ontimer_basket"
sheet_name_3 = "overall_kpis_daywise"
df_1,df_2,df_g,old_text,new_text = slide13.create_df(chogori_ppt_xl_path,sheet_name_1,sheet_name_2,sheet_name_3,overall_abv)
slide_13 = slide13.edit_table(slide_13,df_1,df_2)
slide_13 = slide13.create_overall_program_graph(slide_13,df_g)

for i in range(len(old_text)):
    formatting_functions.change_text_box(slide_13,old_text[i],new_text[i])
#-------------------------------------------------------------------------------------------------------------------------------------------------
# Slide 14
slide_num = 14
slide_14 = prs.slides[slide_num-1]
sheet_name_1 = "avarage_basket_value_mom"
sheet_name_2 = "average_basket_kpis"
df_1,df_2 = slide14.create_df(chogori_ppt_xl_path,sheet_name_1,sheet_name_2)
slide_14 = slide14.edit_table(slide_14,df_1,df_2)
#-------------------------------------------------------------------------------------------------------------------------------------------------
# Slide 15
slide_num = 15
slide_15 = prs.slides[slide_num-1]
sheet_name = "genderwise_kpis"
df = slide15.create_df(chogori_ppt_xl_path,sheet_name)
slide_15 = slide15.edit_table(slide_15,df)
#-------------------------------------------------------------------------------------------------------------------------------------------------
# Slide 16
slide_num = 16
slide_16 = prs.slides[slide_num-1]
sheet_name = "top_10_item_genderwise"
df = slide16.create_df(chogori_ppt_xl_path,sheet_name)
slide_16 = slide16.edit_table(slide_16,df)
#-------------------------------------------------------------------------------------------------------------------------------------------------

#Save the presentation
prs.save(variables.deck_path)

print("Presentation created successfully!")